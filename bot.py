import os
from openai import OpenAI
from telegram import Update
from telegram.ext import ApplicationBuilder, CommandHandler, MessageHandler, filters, ContextTypes
from pptx import Presentation

# 🔐 Render Environment Variables (serverdan o'qiladi)
BOT_TOKEN = os.getenv("8146516889:AAFbvGE3_Hr2xh-j5IvZmaF9E9y-4jzMxlQ")
OPENAI_API_KEY = os.getenv("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")

# 🔍 Tekshiruv (muammo bo‘lsa loglarda chiqadi)
if not BOT_TOKEN:
    raise ValueError("❌ BOT_TOKEN topilmadi. Render Environment Variables da o‘rnatilganini tekshiring.")
if not OPENAI_API_KEY:
    raise ValueError("❌ OPENAI_API_KEY topilmadi. Render Environment Variables da o‘rnatilganini tekshiring.")

client = OpenAI(api_key=OPENAI_API_KEY)

# ✅ /start buyrug‘i
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text("👋 Salom! Menga taqdimot mavzusini yuboring, men sizga PowerPoint fayl tayyorlab beraman 🎓")

# 🧾 Taqdimot yaratish funksiyasi
async def create_presentation(update: Update, context: ContextTypes.DEFAULT_TYPE):
    topic = update.message.text
    await update.message.reply_text("⏳ Taqdimot tayyorlanmoqda...")

    prompt = f"{topic} mavzusi bo‘yicha 5 ta slayddan iborat taqdimot yozing. Har bir slaydda sarlavha va 3-4 ta punkt bo‘lsin."

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": prompt}]
        )
    except Exception as e:
        await update.message.reply_text(f"❌ Xatolik: {e}")
        return

    slides_text = response.choices[0].message.content

    prs = Presentation()
    for slide_text in slides_text.split("Slayd"):
        if slide_text.strip():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            lines = slide_text.strip().split("\n")
            title = lines[0] if lines else "Slayd"
            content = "\n".join(lines[1:]) if len(lines) > 1 else ""
            slide.shapes.title.text = title
            slide.placeholders[1].text = content

    file_name = "presentation.pptx"
    prs.save(file_name)
    await update.message.reply_document(document=open(file_name, "rb"))

# 🧩 Botni ishga tushirish
app = ApplicationBuilder().token(BOT_TOKEN).build()
app.add_handler(CommandHandler("start", start))
app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, create_presentation))

print("✅ Bot ishga tushdi...")
app.run_polling()
